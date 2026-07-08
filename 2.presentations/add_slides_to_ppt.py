"""
Script to add new slides to dap-bpa Introduction PowerPoint presentation
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Slide content data
slides_data = [
    {
        "title": "Introduction to DAP",
        "content": [
            ("What is DAP (Dell Automation Platform)?", [
                "Orchestration Platform: Central platform for automating infrastructure deployments across on-premises and cloud environments",
                "Blueprint-Based: Uses TOSCA-based blueprints to define infrastructure as code",
                "Plugin Architecture: Extensible plugin system supporting multiple platforms (vSphere, AWS, Azure, Kubernetes, etc.)",
                "Multi-Environment: Supports data center, bare-metal, and Kubernetes deployments"
            ]),
            ("Key DAP Components", [
                "Orchestrator: Central engine that executes blueprints and manages deployments",
                "Plugins: Platform-specific extensions providing node types and capabilities",
                "Blueprint Repository: Storage for signed and custom blueprints",
                "Deployment Engine: Manages installation, uninstall, and update workflows",
                "Event Streaming: Real-time monitoring and event tracking"
            ]),
            ("DAP Capabilities", [
                "Deploy complete infrastructure stacks",
                "Manage full deployment lifecycle",
                "Support for intrinsic functions and relationships",
                "Built-in validation and compliance checking",
                "Integration with external systems (SNOW, Datadog, CI/CD)"
            ])
        ]
    },
    {
        "title": "Blueprint Structure",
        "content": [
            ("Blueprint Directory Structure", [
                "my-blueprint/",
                "  ├── blueprint.yaml              # Main blueprint definition",
                "  ├── inputs.yaml                 # Input definitions (optional)",
                "  ├── capabilities.yaml           # Output definitions (optional)",
                "  ├── infrastructure/",
                "  │   └── vsphere/",
                "  │       ├── inputs.yaml",
                "  │       ├── definitions.yaml",
                "  │       └── outputs.yaml",
                "  ├── example_JSON_configs/       # Sample deployment inputs",
                "  │   └── deployment-inputs.json",
                "  ├── CHANGELOG.yaml              # Version history",
                "  ├── README.md                   # Documentation",
                "  └── icon.png                    # Blueprint icon"
            ]),
            ("Core Blueprint Components", [
                "Metadata: name, version, description, author, license",
                "TOSCA Structure: tosca_definitions_version, imports, dsl_definitions",
                "Inputs: User-provided parameters",
                "Node Templates: Infrastructure components",
                "Capabilities: Output values for consumers",
                "Relationships: Dependencies between nodes"
            ]),
            ("Key Sections Explained", [
                "Imports: Load Dell types and required plugins",
                "Inputs: Define parameters users provide at deployment time",
                "Node Templates: Define infrastructure resources using plugin node types",
                "Capabilities: Expose outputs for other blueprints to consume",
                "Relationships: Define dependencies and connections between nodes"
            ])
        ]
    },
    {
        "title": "Introduction to Blueprint Assist (BPA)",
        "content": [
            ("What is Blueprint Assist?", [
                "AI-powered infrastructure deployment tool that simplifies blueprint creation and management for DAP"
            ]),
            ("Key Value Propositions", [
                "Accelerated Development: Reduce blueprint creation time with AI assistance",
                "Consistency: Ensure deployments follow best practices and standards",
                "Knowledge Transfer: Learn from existing blueprints through reasoning and analysis",
                "Flexibility: Support both signed and custom organizational blueprints",
                "Multi-Platform: Deploy across on-premises, bare-metal, and Kubernetes",
                "DAP Integration: Full lifecycle management from CLI or AI chat interface"
            ]),
            ("dap-bpa Architecture", [
                "User → IDE → Skills → LLM → dap-bpa CLI → Blueprint Description → DAP Orchestrator"
            ]),
            ("Core Components", [
                "Skills Engine: AI-powered engine processing natural language requirements",
                "Knowledge Base: Repository of plugins, blueprints, and documentation",
                "dap-bpa CLI: Command-line interface for all operations",
                "IDE Integration: Seamless integration with development environments"
            ]),
            ("Primary Use Cases", [
                "1. Infrastructure Deployment",
                "2. Blueprint Analysis and Auditing",
                "3. Blueprint Creation from Scratch",
                "4. Template Reuse from Signed Blueprints",
                "5. Custom Organization-Specific Blueprints"
            ])
        ]
    },
    {
        "title": "Building Blueprints - Phase 1: Planning",
        "content": [
            ("Define Blueprint Requirements", [
                "Purpose: What infrastructure will this blueprint deploy?",
                "Scope: What components are included/excluded?",
                "Target Environment: On-premises, Kubernetes, cloud?",
                "Constraints: Budget, timeline, compliance requirements?",
                "Dependencies: External systems or services required?"
            ]),
            ("Research Existing Blueprints", [
                "dap-bpa knowledge blueprints find \"storage\"",
                "dap-bpa knowledge blueprints find \"network\"",
                "dap-bpa knowledge blueprints get vsphere-vm --include-files"
            ]),
            ("Study Plugin Documentation", [
                "dap-bpa knowledge plugins list vsphere",
                "dap-bpa knowledge plugins get vsphere dell.nodes.vsphere.Server",
                "dap-bpa knowledge plugins docs vsphere"
            ]),
            ("Create Blueprint Directory", [
                "mkdir my-blueprint",
                "cd my-blueprint",
                "mkdir -p infrastructure/vsphere",
                "mkdir -p example_JSON_configs"
            ])
        ]
    },
    {
        "title": "Building Blueprints - Phase 2: Authoring",
        "content": [
            ("Blueprint File Structure Options", [
                "Simple Blueprint (single file): blueprint.yaml, example configs, changelog, readme, icon",
                "Complex Blueprint (multi-file): Main blueprint with imports, separate inputs/capabilities/definitions files"
            ]),
            ("Key TOSCA Elements - Imports", [
                "imports:",
                "  - dell/types/types.yaml",
                "  - plugin:vsphere-plugin?version= >=3.0.7.0,<4.0.0.0"
            ]),
            ("DSL Definitions (Reusable Config)", [
                "dsl_definitions:",
                "  connection_config: &connection_config",
                "    username: { get_input: vsphere_username }",
                "    password: { get_secret: { get_input: vsphere_secret } }"
            ]),
            ("Node Templates", [
                "node_templates:",
                "  my_vm:",
                "    type: dell.nodes.vsphere.Server",
                "    properties:",
                "      connection_config: *connection_config"
            ]),
            ("Intrinsic Functions", [
                "{ get_input: parameter_name }",
                "{ get_secret: secret_name }",
                "{ get_attribute: [node_name, attribute] }",
                "{ concat: [\"prefix-\", { get_input: name }] }"
            ])
        ]
    },
    {
        "title": "Building Blueprints - Phase 3: Validation & Deployment",
        "content": [
            ("Validation Steps", [
                "1. Local Linting: dap-bpa blueprint lint --file blueprint.yaml",
                "2. Node Validation: dap-bpa blueprint validate my_vm --file blueprint.yaml",
                "3. Create Example Configurations: JSON with deployment inputs"
            ]),
            ("Deployment Steps", [
                "1. Upload Blueprint:",
                "   dap-bpa orchestrator blueprints upload --file blueprint.yaml --id my-blueprint --revision 1.0.0",
                "2. Create Deployment:",
                "   dap-bpa orchestrator deployments create --blueprint-id my-blueprint --inputs example_JSON_configs/deployment-inputs.json",
                "3. Execute Installation:",
                "   dap-bpa orchestrator executions start <deployment_id> --workflow-id install"
            ]),
            ("Monitoring", [
                "dap-bpa orchestrator deployments get <deployment_id>",
                "dap-bpa orchestrator events get <execution_id>"
            ])
        ]
    },
    {
        "title": "Adding Skills to Blueprints",
        "content": [
            ("What are dap-bpa Skills?", [
                "Skills map directly to DAP plugins and operational workflows",
                "Provide node type knowledge, authoring guidance, documentation lookup",
                "Natural-language activation in IDE"
            ]),
            ("Skill Categories", [
                "Plugin Skills: dap-plugin-aws, dap-plugin-azure, dap-plugin-vsphere, dap-plugin-kubernetes, etc.",
                "Workflow Skills: dap-deployment-update, dap-service-composition, dap-monitor",
                "Authoring Skills: dap-blueprint-assist, dap-scripts, dap-sdd-guidelines"
            ]),
            ("Using Skills in IDE", [
                "Skills are invoked through natural-language prompts:",
                "\"Help me create a vSphere VM blueprint\"",
                "\"Add Kubernetes deployment to my blueprint\"",
                "\"Show me the AWS EC2 node type properties\""
            ]),
            ("Querying Skills via CLI", [
                "dap-bpa skills list --detailed",
                "dap-bpa knowledge plugins list kubernetes",
                "dap-bpa knowledge plugins get kubernetes dell.nodes.kubernetes.Deployment",
                "dap-bpa knowledge docs search \"helm chart deployment\"",
                "dap-bpa knowledge blueprints find \"kubernetes helm\""
            ])
        ]
    },
    {
        "title": "Skills Architecture",
        "content": [
            ("Skills-Based Architecture Principles", [
                "1. Modularity: Each skill encapsulates a specific capability",
                "2. Composability: Skills can be combined to create complex deployments",
                "3. Reusability: Skills can be reused across multiple blueprints",
                "4. Independence: Skills can be developed, tested, and updated independently",
                "5. Discoverability: Skills are easily discoverable and self-documenting"
            ]),
            ("Skill Composition Patterns", [
                "Linear: Skill A → Skill B → Skill C → Skill D",
                "Parallel: [Network, Storage] → Compute",
                "Conditional: Skill A → (Condition?) → Skill B : Skill C"
            ]),
            ("Skill Design Principles", [
                "Single Responsibility: Each skill has one clear purpose",
                "Clear Interfaces: Well-defined input/output contracts",
                "Fail Fast: Validate inputs early",
                "Idempotency: Consistent results on re-execution",
                "Observability: Logging and progress reporting"
            ]),
            ("Skill Files Location", [
                "~/.blueprint-assist/skills/"
            ])
        ]
    },
    {
        "title": "Best Practices Summary",
        "content": [
            ("Blueprint Development", [
                "Use clear, descriptive names for node templates",
                "Group related inputs using input_groups",
                "Use dsl_definitions for reusable values",
                "Separate complex blueprints into multiple files",
                "Always lint before uploading",
                "Validate node templates individually"
            ]),
            ("Input Management", [
                "Provide default values where appropriate",
                "Use display groups for better UX",
                "Include helpful descriptions for all inputs",
                "Use proper constraints for validation",
                "Separate configuration from code"
            ]),
            ("Documentation", [
                "Maintain comprehensive README files",
                "Keep CHANGELOG.yaml updated",
                "Document all parameters and their purposes",
                "Include troubleshooting guides",
                "Create example configurations"
            ]),
            ("Validation & Testing", [
                "Always lint before uploading",
                "Validate node templates individually",
                "Test with example configurations",
                "Use monitor for end-to-end testing",
                "Test in multiple environments"
            ]),
            ("Skills Usage", [
                "Leverage existing skills before creating custom ones",
                "Use natural-language prompts in IDE",
                "Query knowledge base for node types and examples",
                "Study plugin documentation before authoring",
                "Use signed blueprints as starting points"
            ])
        ]
    },
    {
        "title": "Resources & Next Steps",
        "content": [
            ("Key Commands Reference", [
                "Blueprint Operations: dap-bpa blueprint lint, dap-bpa orchestrator blueprints upload",
                "Knowledge Base: dap-bpa knowledge blueprints find, dap-bpa knowledge plugins list",
                "Deployment Operations: dap-bpa orchestrator deployments create, dap-bpa orchestrator executions start"
            ]),
            ("Repository Sections", [
                "Section 001: Introduction to Blueprint Assist",
                "Section 004: Skills Overview",
                "Section 005: Skills Architecture",
                "Section 006: Supported Blueprints",
                "Section 007: Building Blueprints",
                "Section 010: Blueprint Anatomy",
                "Section 011: Skill Anatomy",
                "Section 013: dap-bpa CLI Commands"
            ]),
            ("Internal Blueprints Repository", [
                "ISG-Edge/hzp-eo-solutions",
                "Kubernetes on Bare-metal, CSI PowerStore",
                "PowerStore Storage, Bare Metal OS Deployment"
            ]),
            ("Next Steps", [
                "1. Complete hands-on workshop (Section 012)",
                "2. Practice with signed blueprints",
                "3. Build your first custom blueprint",
                "4. Explore advanced features",
                "5. Contribute to internal repository"
            ])
        ]
    }
]

def add_slide_with_content(prs, title, content_sections):
    """Add a slide with title and bullet point content sections"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Get content placeholder
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    # Add content sections
    for section_title, bullets in content_sections:
        # Add section title
        p = text_frame.add_paragraph()
        p.text = section_title
        p.font.bold = True
        p.font.size = Pt(18)
        p.space_after = Pt(6)
        
        # Add bullets
        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 1
            p.font.size = Pt(14)
            p.space_after = Pt(3)
        
        # Add spacing between sections
        if content_sections.index((section_title, bullets)) < len(content_sections) - 1:
            p = text_frame.add_paragraph()
            p.space_after = Pt(10)

def main():
    # Load existing presentation
    ppt_path = r"1.BPA-intro.pptx"
    
    try:
        prs = Presentation(ppt_path)
        print(f"Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"Error loading presentation: {e}")
        return
    
    # Add new slides
    for slide_data in slides_data:
        try:
            add_slide_with_content(prs, slide_data["title"], slide_data["content"])
            print(f"Added slide: {slide_data['title']}")
        except Exception as e:
            print(f"Error adding slide '{slide_data['title']}': {e}")
    
    # Save the presentation
    output_path = r"1.BPA-intro-updated.pptx"
    try:
        prs.save(output_path)
        print(f"\nSaved updated presentation to: {output_path}")
        print(f"Total slides: {len(prs.slides)}")
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == "__main__":
    main()
