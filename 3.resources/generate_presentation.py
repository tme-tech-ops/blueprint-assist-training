#!/usr/bin/env python3
"""
Generate PowerPoint presentation from template using section content files.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def read_markdown_sections():
    """Read all section markdown files and parse slide content."""
    sections = []
    sections_dir = "sections"
    
    section_dirs = [
        "section-1-introduction",
        "section-2-installation",
        "section-3-skills-overview",
        "section-4-skills-architecture",
        "section-5-building-blueprints",
        "section-6-blueprint-reasoning",
        "section-7-supported-blueprints",
        "section-8-blueprint-anatomy",
        "section-9-skill-anatomy",
        "section-10-hands-on-workshop"
    ]
    
    for section_dir in section_dirs:
        file_path = os.path.join(sections_dir, section_dir, "content.md")
        if os.path.exists(file_path):
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                slides = parse_slides(content)
                sections.extend(slides)
    
    return sections

def parse_slides(markdown_content):
    """Parse markdown content into slide data."""
    slides = []
    lines = markdown_content.split('\n')
    
    current_title = None
    current_bullets = []
    
    for line in lines:
        line = line.strip()
        
        # Section header (#) - skip these
        if line.startswith('# ') and not line.startswith('##'):
            continue
        
        # Topic header (##) - create new slide
        elif line.startswith('## '):
            # Save previous slide if exists
            if current_title and current_bullets:
                slides.append({'title': current_title, 'bullets': current_bullets})
            
            # Start new slide
            current_title = line.replace('## ', '')
            current_bullets = []
        
        # Sub-topic header (###) - add as bullet point
        elif line.startswith('### '):
            current_bullets.append(line.replace('### ', ''))
        
        # Bullet points
        elif line.startswith('- '):
            current_bullets.append(line[2:])
        elif line.startswith('* '):
            current_bullets.append(line[2:])
        elif line.startswith('1. ') or line.startswith('2. ') or line.startswith('3. ') or line.startswith('4. ') or line.startswith('5. '):
            current_bullets.append(line[3:])
        elif line.startswith('- **'):
            # Bold bullet points
            current_bullets.append(line[2:])
        elif line and current_title and not line.startswith('#'):
            # Regular text content
            if len(line) < 200:  # Only add short lines as bullets
                current_bullets.append(line)
    
    # Add last slide
    if current_title and current_bullets:
        slides.append({'title': current_title, 'bullets': current_bullets})
    
    return slides

def create_presentation(slides, template_path, output_path):
    """Create PowerPoint presentation from slides data."""
    # Load template
    prs = Presentation(template_path)
    
    # Clear existing slides (keep first slide if it's a title slide)
    for i in range(len(prs.slides) - 1, 0, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    # Add new slides
    for slide_data in slides:
        add_slide(prs, slide_data)
    
    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

def add_slide(prs, slide_data):
    """Add a slide with title and bullet points."""
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1.2)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 42, 84)  # Dark blue
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add bullet points
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9)
    height = Inches(5.5)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    
    for i, bullet in enumerate(slide_data['bullets']):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
        p.space_after = Pt(14)
        p.line_spacing = 1.3
        
        # Check if it's a sub-bullet (indented)
        if bullet.startswith('  ') or bullet.startswith('\t'):
            p.level = 1
            p.font.size = Pt(18)
            p.space_before = Pt(8)

def main():
    """Main function to generate presentation."""
    print("Reading markdown sections...")
    slides = read_markdown_sections()
    print(f"Found {len(slides)} slides")
    
    template_path = "presentations/BPA-intro-template.pptx"
    output_path = "presentations/BPA-EKT-Training-Session.pptx"
    
    print(f"Loading template: {template_path}")
    print(f"Generating presentation: {output_path}")
    
    create_presentation(slides, template_path, output_path)
    
    print("Done!")

if __name__ == "__main__":
    main()
